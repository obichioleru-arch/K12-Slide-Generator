"""Safe editable PPTX exporter for the District Slide Tool.

This version intentionally avoids native PowerPoint chart objects because those
can create fragile chart XML across PowerPoint versions. Instead, charts are
drawn with editable PowerPoint rectangles, lines, and text boxes. The output is
still editable in PowerPoint, but it opens without PowerPoint repair prompts.
"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


NAVY = "003291"
RED = "D97706"
CYAN = "00B0F0"
WHITE = "FFFFFF"
LIGHT_BLUE = "EAF6FF"
LIGHT_GREEN = "ECFDF5"
LIGHT_GRAY = "F3F4F6"
GRID_GRAY = "E5E7EB"
MID_GRAY = "9CA3AF"
DARK_TEXT = "1F2933"
GREEN = "16A34A"
AMBER = "D97706"
GOLD = "FFC000"
PURPLE = "C084FC"

COLORS = [NAVY, CYAN, RED, GOLD, PURPLE, "6B7280"]


def _series_color(name: Any, idx: int = 0) -> str:
    n = str(name or "").lower()
    if "not met" in n:
        return RED
    if "tsi met" in n or n.strip() == "met" or " met" in n:
        return NAVY
    return COLORS[idx % len(COLORS)]




def _clean_source_text(text: str) -> str:
    s = str(text or "").strip()
    s = re.sub(r"^(Source:\s*[^.]+\.?\s*)+(?=Source:\s*)", "", s, flags=re.I)
    s = re.sub(r"^Source:\s*District Salesforce\.\s*(?=Source:\s*)", "", s, flags=re.I)
    return s.strip()

def _rgb(hex_color: str) -> RGBColor:
    h = str(hex_color).replace("#", "")
    if len(h) != 6:
        h = NAVY
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _emu(px: float):
    """Convert from the app's 1333 x 750 pixel slide system to pptx inches."""
    return Inches(float(px) / 100.0)


def _add_rect(slide, x, y, w, h, fill, line=None, radius=False, alpha=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, _emu(x), _emu(y), _emu(w), _emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    shp.line.color.rgb = _rgb(line or fill)
    shp.line.width = Pt(0.5)
    return shp


def _add_oval(slide, x, y, w, h, fill=None, line=NAVY, width=2):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, _emu(x), _emu(y), _emu(w), _emu(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    else:
        shp.fill.background()
    shp.line.color.rgb = _rgb(line)
    shp.line.width = Pt(width)
    return shp


def _add_line(slide, x1, y1, x2, y2, color=GRID_GRAY, width=1.0):
    line = slide.shapes.add_connector(1, _emu(x1), _emu(y1), _emu(x2), _emu(y2))
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = DARK_TEXT,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return box


def _fmt(v: Any, mode: str = "count") -> str:
    try:
        n = float(v)
    except Exception:
        return str(v)
    if mode == "percent":
        return f"{n:.1f}%"
    if abs(n) >= 1_000_000:
        return f"${n/1_000_000:.1f}M" if n > 100000 else f"{n/1_000_000:.1f}M"
    if abs(n) >= 1000:
        return f"{n:,.0f}"
    return f"{n:.0f}" if abs(n - round(n)) < 0.001 else f"{n:.1f}"


def _num(v: Any) -> float:
    try:
        if isinstance(v, str):
            return float(v.replace("%", "").replace(",", "").replace("$", ""))
        return float(v)
    except Exception:
        return 0.0


def _short(s: Any, max_len: int = 22) -> str:
    s = str(s or "")
    return s if len(s) <= max_len else s[: max_len - 1] + "…"


def _fmt_money(v: Any) -> str:
    try:
        n = float(str(v).replace("$", "").replace(",", ""))
    except Exception:
        return "$0"
    if abs(n) >= 1_000_000:
        return f"${n/1_000_000:.1f}M"
    if abs(n) >= 1_000:
        return f"${n/1_000:.1f}K"
    if abs(n) >= 1:
        # Some calculator payloads may already store HB3 values in millions.
        return f"${n:.1f}M"
    if n > 0:
        return f"${n*1000:.0f}K"
    return "$0"

def _hb3_status_for(cat: Any, fallback: str = "verified") -> str:
    m = re.search(r"(20\d{2})", str(cat or ""))
    yr = int(m.group(1)) if m else None
    if yr == 2025:
        return "estimate"
    if yr == 2026:
        return "projected"
    return fallback or "verified"




def _campus_label(name: Any) -> str:
    s = str(name or "").strip()
    replacements = [
        ("High School", "HS"),
        ("Accelerated", "Accel."),
        ("Academy", "Acad."),
        ("Collegiate Institute", "Collegiate Inst."),
        ("Career High", "Career"),
        ("Young Women's Leadership Academy", "YWLA"),
        ("at Bill Arnold", "Bill Arnold"),
    ]
    for old, new in replacements:
        s = s.replace(old, new)
    return s

def _find_logo_path() -> str | None:
    """Find the PPTX footer logo only.

    This intentionally does NOT use frontend/src/assets/emc_landing_logo.png,
    because that file is reserved for the landing page.
    """
    here = Path(__file__).resolve()
    candidates = [
        here.parent / "emc_footer_logo.png",
        Path.cwd() / "backend" / "emc_footer_logo.png",
        here.parent / "emc_logo.png",
        Path.cwd() / "backend" / "emc_logo.png",
    ]
    for p in candidates:
        if p.exists() and p.is_file():
            return str(p)
    return None



def _asset_path(name: str) -> str | None:
    here = Path(__file__).resolve().parent
    candidates = [here / name, Path.cwd() / "backend" / name, Path.cwd() / name]
    for p in candidates:
        if p.exists() and p.is_file():
            return str(p)
    return None


def _add_picture_asset(slide, name: str, x: float, y: float, w: float, h: float | None = None) -> bool:
    path = _asset_path(name)
    if not path:
        return False
    try:
        if h is None:
            slide.shapes.add_picture(path, _emu(x), _emu(y), width=_emu(w))
        else:
            slide.shapes.add_picture(path, _emu(x), _emu(y), width=_emu(w), height=_emu(h))
        return True
    except Exception:
        return False


def _add_dark_background_logo(slide, x: float, y: float, w: float, h: float | None = None) -> bool:
    """White-text EMC logo for dark blue cover/mission slides."""
    if _add_picture_asset(slide, "emc_logo_white_transparent.png", x, y, w, h):
        return True
    return _add_logo(slide, x, y, w, h)

def _add_logo(slide, x: float, y: float, w: float, h: float | None = None):
    """Add the real EMC logo when available; fallback to text only."""
    logo = _find_logo_path()
    if logo:
        try:
            if h is None:
                slide.shapes.add_picture(logo, _emu(x), _emu(y), width=_emu(w))
            else:
                slide.shapes.add_picture(logo, _emu(x), _emu(y), width=_emu(w), height=_emu(h))
            return True
        except Exception:
            pass
    _add_text(slide, "ECONOMIC MOBILITY\nCENTER", x, y + (h or 30), w, 24, size=6.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    return False


def _clean_title(sc: Dict[str, Any]) -> str:
    sd = sc.get("slide_data", {}) or {}
    if sd.get("Title"):
        return str(sd.get("Title"))
    return str(sc.get("slide_type", "Slide")).replace("_", " ").title()


def _title(slide, title: str, district: str = "", month: str = "", year: str = "", tag: str = "DATA"):
    # Header layout mirrors the HTML template but keeps district and date badge separate.
    _add_text(slide, title, 0, 24, 1333, 34, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    if district:
        _add_text(slide, district, 0, 64, 1333, 18, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    stamp = " ".join([x for x in [month, year] if x])
    if stamp:
        # Date badge intentionally sits below the district label to avoid overlap.
        _add_rect(slide, 616, 86, 104, 22, NAVY, radius=True)
        _add_text(slide, stamp, 620, 91, 96, 12, size=8, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    _add_line(slide, 28, 124, 1305, 124, NAVY, width=1.6)
    _add_line(slide, 28, 126, 1305, 126, CYAN, width=0.8)

    if tag:
        _add_rect(slide, 1248, 18, 55, 24, RED, radius=True)
        _add_text(slide, tag, 1252, 23, 47, 12, size=8, color=WHITE, bold=True, italic=True, align=PP_ALIGN.CENTER)


def _footer(slide, footnote: str = ""):
    footnote = _clean_source_text(footnote)
    # Use the PPTX-specific EMC footer logo. This asset is separate from
    # the frontend landing-page logo so the landing page cannot be changed
    # by PowerPoint export work.
    _add_logo(slide, 42, 682, 82)

    if footnote:
        _add_text(slide, footnote, 142, 717, 960, 16, size=7, color="4B5563")

    _add_rect(slide, 1224, 724, 80, 3, RED)


def _add_insight_boxes(slide, insights: Sequence[str], x=32, y=585, w=1270, h=92):
    half = (w - 14) / 2
    _add_rect(slide, x, y, half, h, LIGHT_BLUE, line=LIGHT_BLUE, radius=True)
    _add_rect(slide, x, y, 4, h, CYAN)
    _add_text(slide, "KEY INSIGHTS", x + 16, y + 12, half - 30, 18, size=12, color=NAVY, bold=True)
    insight_lines = list(insights or [])
    for i, text in enumerate(insight_lines[:2]):
        _add_text(slide, "▶", x + 18, y + 36 + i * 25, 14, 12, size=9, color=NAVY, bold=True)
        _add_text(slide, str(text).lstrip("•▶- ").strip(), x + 34, y + 35 + i * 25, half - 48, 20, size=7.5, color=DARK_TEXT)
        if i == 0:
            _add_line(slide, x + 18, y + 57, x + half - 14, y + 57, CYAN, width=0.35)

    nx = x + half + 14
    _add_rect(slide, nx, y, half, h, LIGHT_GREEN, line=LIGHT_GREEN, radius=True)
    _add_rect(slide, nx, y, 4, h, GREEN)
    _add_text(slide, "NEXT STEPS", nx + 16, y + 12, half - 30, 18, size=12, color=GREEN, bold=True)
    next_text = insight_lines[2] if len(insight_lines) >= 3 else "Review the data and identify the highest-priority follow-up action."
    _add_text(slide, "▶", nx + 18, y + 40, 14, 12, size=9, color=GREEN, bold=True)
    _add_text(slide, str(next_text).lstrip("•▶- ").strip(), nx + 34, y + 38, half - 48, 42, size=7.5, color="065F46")
    _add_line(slide, nx + 18, y + 71, nx + half - 14, y + 71, GREEN, width=0.35)


def _legend(slide, series: Sequence[Dict[str, Any]], x=100, y=548, max_w=900):
    cx = x
    for idx, s in enumerate(series[:6]):
        name = str(s.get("name", f"Series {idx+1}"))
        _add_rect(slide, cx, y, 14, 10, _series_color(s.get("name", ""), idx), radius=True)
        _add_text(slide, name, cx + 20, y - 2, 150, 15, size=8, color=DARK_TEXT)
        cx += min(210, max(110, len(name) * 7 + 50))
        if cx > x + max_w:
            break


def _draw_stacked_columns(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    if not categories or not series:
        _add_text(slide, "No chart data available", 100, 260, 900, 50, size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return

    x0, y0, w, h = 115, 130, 1165, 398
    _add_line(slide, x0, y0, x0, y0 + h, MID_GRAY)
    _add_line(slide, x0, y0 + h, x0 + w, y0 + h, MID_GRAY)

    totals = []
    for ci in range(len(categories)):
        totals.append(sum(_num(s.get("values", [0] * len(categories))[ci]) for s in series))
    max_v = max(totals + [100 if mode == "percent" else 1])
    max_v = max(1, max_v)
    if mode == "percent":
        max_v = max(100, max_v)

    tick_count = 5
    for i in range(tick_count + 1):
        val = max_v * i / tick_count
        y = y0 + h - (val / max_v) * h
        _add_line(slide, x0, y, x0 + w, y, GRID_GRAY, width=0.5)
        _add_text(slide, _fmt(val, mode), x0 - 55, y - 7, 44, 12, size=7, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    gap = w / max(1, len(categories))
    bar_w = min(70, gap * 0.36)
    for ci, cat in enumerate(categories):
        cx = x0 + gap * ci + gap / 2
        y_cursor = y0 + h
        for si, s in enumerate(series):
            vals = s.get("values", [])
            val = _num(vals[ci]) if ci < len(vals) else 0
            bh = (val / max_v) * h
            y_cursor -= bh
            _add_rect(slide, cx - bar_w / 2, y_cursor, bar_w, max(1, bh), _series_color(s.get("name", ""), si), radius=True)
            if bh > 22:
                _add_text(slide, _fmt(val, mode), cx - bar_w / 2, y_cursor + bh / 2 - 6, bar_w, 12, size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, _short(cat, 18), cx - 60, y0 + h + 16, 120, 28, size=7, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    axis_label = "%" if mode == "percent" else "# Students"
    _add_text(slide, axis_label, 40, 300, 45, 18, size=8, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    _legend(slide, series, x=120, y=555)


def _draw_grouped_columns(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    if not categories or not series:
        _add_text(slide, "No chart data available", 100, 260, 900, 50, size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return

    x0, y0, w, h = 115, 130, 1165, 398
    max_v = max([_num(v) for s in series for v in s.get("values", [])] + [100 if mode == "percent" else 1])
    if mode == "percent":
        max_v = max(100, max_v)
    max_v = max(1, max_v)

    _add_line(slide, x0, y0, x0, y0 + h, MID_GRAY)
    _add_line(slide, x0, y0 + h, x0 + w, y0 + h, MID_GRAY)
    for i in range(6):
        val = max_v * i / 5
        y = y0 + h - (val / max_v) * h
        _add_line(slide, x0, y, x0 + w, y, GRID_GRAY, width=0.5)
        _add_text(slide, _fmt(val, mode), x0 - 55, y - 7, 44, 12, size=7, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    group_gap = w / max(1, len(categories))
    bar_w = min(36, (group_gap * 0.72) / max(1, len(series)))
    for ci, cat in enumerate(categories):
        group_left = x0 + group_gap * ci + group_gap / 2 - (bar_w * len(series)) / 2
        for si, s in enumerate(series):
            vals = s.get("values", [])
            val = _num(vals[ci]) if ci < len(vals) else 0
            bh = (val / max_v) * h
            x = group_left + si * bar_w
            y = y0 + h - bh
            _add_rect(slide, x, y, bar_w * 0.84, max(1, bh), _series_color(s.get("name", ""), si), radius=True)
            if bh > 25:
                _add_text(slide, _fmt(val, mode), x, y + bh / 2 - 6, bar_w * 0.84, 12, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, _short(cat, 18), x0 + group_gap * ci + 5, y0 + h + 16, group_gap - 10, 26, size=7, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    _legend(slide, series, x=120, y=555)


def _draw_horizontal_bars(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    if not categories or not series:
        _add_text(slide, "No chart data available", 100, 260, 900, 50, size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return

    x0, y0, w, h = 420, 145, 720, 380
    vals = series[0].get("values", [])
    max_v = max([_num(v) for v in vals] + [100 if mode == "percent" else 1])
    if mode == "percent":
        max_v = max(100, max_v)
    max_v = max(1, max_v)

    for i in range(6):
        val = max_v * i / 5
        x = x0 + (val / max_v) * w
        _add_line(slide, x, y0, x, y0 + h, GRID_GRAY, width=0.5)
        _add_text(slide, _fmt(val, mode), x - 22, y0 + h + 12, 44, 12, size=7, color=MID_GRAY, align=PP_ALIGN.CENTER)

    row_gap = h / max(1, len(categories))
    bar_h = min(34, row_gap * 0.58)
    for ci, cat in enumerate(categories):
        val = _num(vals[ci]) if ci < len(vals) else 0
        bw = (val / max_v) * w
        y = y0 + row_gap * ci + (row_gap - bar_h) / 2
        _add_text(slide, _campus_label(cat), 60, y + 3, 345, 16, size=9, color=DARK_TEXT, align=PP_ALIGN.RIGHT)
        _add_rect(slide, x0, y, max(2, bw), bar_h, COLORS[0], radius=True)
        if bw > 60:
            _add_text(slide, _fmt(val, mode), x0 + bw / 2 - 35, y + bar_h / 2 - 7, 70, 14, size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, _fmt(val, mode), x0 + bw + 8, y + bar_h / 2 - 7, 70, 14, size=8, color=DARK_TEXT, bold=True)
    _add_line(slide, x0, y0, x0, y0 + h, MID_GRAY)
    _add_line(slide, x0, y0 + h, x0 + w, y0 + h, MID_GRAY)



def _draw_horizontal_stacked_bars(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    if not categories or not series:
        _add_text(slide, "No chart data available", 100, 260, 900, 50, size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return

    x0, y0, w, h = 420, 145, 720, 380
    max_v = 100 if mode == "percent" else max(
        [sum(_num(s.get("values", [0] * len(categories))[ci]) for s in series) for ci in range(len(categories))] + [1]
    )

    for i in range(6):
        val = max_v * i / 5
        x = x0 + (val / max_v) * w
        _add_line(slide, x, y0, x, y0 + h, GRID_GRAY, width=0.5)
        _add_text(slide, _fmt(val, mode), x - 22, y0 + h + 12, 44, 12, size=7, color=MID_GRAY, align=PP_ALIGN.CENTER)

    row_gap = h / max(1, len(categories))
    bar_h = min(34, row_gap * 0.58)

    for ci, cat in enumerate(categories):
        y = y0 + row_gap * ci + (row_gap - bar_h) / 2
        _add_text(slide, _campus_label(cat), 60, y + 3, 345, 16, size=9, color=DARK_TEXT, align=PP_ALIGN.RIGHT)
        x_cursor = x0
        for si, s in enumerate(series):
            vals = s.get("values", [])
            val = _num(vals[ci]) if ci < len(vals) else 0
            bw = (val / max_v) * w if max_v else 0
            _add_rect(slide, x_cursor, y, max(2, bw), bar_h, _series_color(s.get("name", ""), si), radius=True)
            if bw > 48:
                _add_text(slide, _fmt(val, mode), x_cursor + bw / 2 - 28, y + bar_h / 2 - 7, 56, 14, size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            x_cursor += bw

    _add_line(slide, x0, y0, x0, y0 + h, MID_GRAY)
    _add_line(slide, x0, y0 + h, x0 + w, y0 + h, MID_GRAY)
    _legend(slide, series, x=360, y=555)

def _draw_card_status(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    values = []
    if (chart_data or {}).get("series"):
        values = (chart_data["series"][0] or {}).get("values", [])
    colors = [GREEN, AMBER, RED]
    x = 310
    for i, cat in enumerate(categories[:3]):
        val = values[i] if i < len(values) else 0
        _add_rect(slide, x + i * 300, 128, 250, 96, ["ECFDF5", "FFFBEB", "FEF2F2"][i % 3], line=colors[i % 3], radius=True)
        _add_text(slide, cat.upper(), x + i * 300 + 15, 146, 220, 18, size=13, color=colors[i % 3], bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, _fmt(val, mode), x + i * 300 + 15, 170, 220, 34, size=26, color=colors[i % 3], bold=True, align=PP_ALIGN.CENTER)


def _draw_by_numbers(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    vals = series[0].get("values", []) if series else []
    positions = [(260, 200), (560, 200), (860, 200)]
    for i, (x, y) in enumerate(positions):
        label = categories[i] if i < len(categories) else f"Metric {i+1}"
        val = vals[i] if i < len(vals) else 0
        _add_oval(slide, x, y, 210, 210, fill=NAVY, line=CYAN, width=4)
        _add_text(slide, _fmt(val, mode), x, y + 70, 210, 50, size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, label, x - 40, y + 240, 290, 28, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)









def _draw_ccmr_pathway_dashboard(slide, chart_data: Dict[str, Any], sc: Dict[str, Any]):
    sd = sc.get("slide_data", {}) or {}
    total = int(_num(chart_data.get("total_students") or sd.get("total_students") or 0))
    on_pathway = int(_num(chart_data.get("on_pathway") or sd.get("on_pathway") or 0))
    on_pct = _num(chart_data.get("on_pathway_pct") if chart_data.get("on_pathway_pct") is not None else sd.get("on_pathway_pct", 0))
    off_pathway = int(_num(chart_data.get("not_on_pathway") or sd.get("not_on_pathway") or 0))
    off_pct = _num(chart_data.get("not_on_pathway_pct") if chart_data.get("not_on_pathway_pct") is not None else sd.get("not_on_pathway_pct", 0))

    counts = chart_data.get("counts") or []
    pcts = chart_data.get("percentages") or []
    vals = chart_data.get("series", [{}])[0].get("values", [])
    mode = str(chart_data.get("mode") or sc.get("mode") or "count").lower()
    if not counts:
        if mode == "percent":
            pcts = vals
            counts = [round(_num(p) / 100 * max(total,1)) for p in pcts]
        else:
            counts = vals
            pcts = [round(_num(v) / max(total,1) * 100, 1) for v in counts]
    if not pcts:
        pcts = [round(_num(v) / max(total,1) * 100, 1) for v in counts]

    district = sd.get("District") or ""
    month = str(sc.get("month","")).strip()
    year = str(sc.get("year_label", sc.get("year",""))).strip()
    date_line = f"{month} {year}".strip()

    _add_line(slide, 0, 0, 1333, 0, NAVY, width=2)
    _add_text(slide, "CCMR PATHWAY ANALYSIS", 28, 28, 400, 18, size=13.5, color="E8192C", bold=True)
    _add_text(slide, str(district), 28, 56, 420, 25, size=18, color=NAVY, bold=True)
    _add_text(slide, f"{date_line} • Total Enrollment: {total:,} students", 28, 88, 600, 18, size=12.5, color=DARK_TEXT)
    _add_line(slide, 28, 118, 1305, 118, CYAN, width=1.2)
    _add_rect(slide, 1210, 22, 90, 24, "E8192C", radius=True)
    _add_text(slide, "CCMR Pathways", 1210, 29, 90, 9, size=6.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _add_rect(slide, 28, 146, 222, 360, "E8192C", line="E8192C", radius=True)
    _add_text(slide, "!", 119, 170, 40, 28, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"{off_pathway:,}", 28, 238, 222, 46, size=31, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Students", 28, 300, 222, 18, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "NOT ON ANY\nCCMR PATHWAY", 28, 332, 222, 45, size=12.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, 50, 442, 178, 23, "F95B68", line="F95B68", radius=True)
    _add_text(slide, f"{off_pct:.1f}% of all students", 50, 449, 178, 9, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _add_rect(slide, 267, 146, 1038, 70, NAVY, line=NAVY, radius=True)
    _add_text(slide, "STUDENTS ON A CCMR PATHWAY", 285, 164, 500, 15, size=12, color="BFD7FF", bold=True)
    _add_text(slide, f"{on_pathway:,} students are on pathway toward CCMR", 285, 193, 390, 12, size=11, color=WHITE, bold=True)
    _add_text(slide, f"{on_pct:.1f}% of enrollment", 650, 193, 250, 12, size=11, color=CYAN, bold=True)

    labels = ["DUAL CREDIT", "TSI", "AP/IB", "IBC"]
    colors = [CYAN, "7C3AED", "10B981", "F59E0B"]
    card_x, card_y, card_w, card_h, gap = 267, 232, 250, 110, 8
    for i, label in enumerate(labels):
        x = card_x + i * (card_w + gap)
        count = int(_num(counts[i])) if i < len(counts) else 0
        pct = _num(pcts[i]) if i < len(pcts) else 0
        _add_rect(slide, x, card_y, card_w, card_h, "F3F4F6", line="F3F4F6", radius=True)
        _add_line(slide, x, card_y, x + card_w, card_y, colors[i], width=2)
        _add_text(slide, label, x, card_y + 20, card_w, 12, size=10.5, color=colors[i], bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"{count:,}", x, card_y + 48, card_w, 22, size=20, color="111827", bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"{pct:.1f}%", x, card_y + 78, card_w, 13, size=10.5, color=colors[i], bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, "students", x, card_y + 94, card_w, 9, size=6.8, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    top_idx = max(range(min(len(pcts),4)), key=lambda k: _num(pcts[k])) if pcts else 0
    top_label = labels[top_idx]
    top_count = int(_num(counts[top_idx])) if top_idx < len(counts) else 0
    top_pct = _num(pcts[top_idx]) if top_idx < len(pcts) else 0

    _add_rect(slide, 267, 358, 1038, 150, NAVY, line=NAVY, radius=True)
    _add_text(slide, "ACTION INSIGHTS", 285, 376, 360, 16, size=12, color=CYAN, bold=True)
    _add_text(slide, f"▶ {on_pathway:,} students ({on_pct:.1f}%) are on a CCMR pathway — {off_pathway:,} ({off_pct:.1f}%) have no pathway and need immediate outreach.", 285, 405, 980, 16, size=10.2, color=WHITE)
    _add_text(slide, f"▶ {top_label} is the most common pathway with {top_count:,} students ({top_pct:.1f}% of enrollment).", 285, 434, 980, 16, size=10.2, color=WHITE)
    _add_text(slide, f"▶ Targeted intervention for the {off_pathway:,} off-pathway students could significantly improve CCMR outcomes.", 285, 463, 980, 16, size=10.2, color=WHITE)

def _draw_ccmr_af_status_dashboard(slide, chart_data: Dict[str, Any], sc: Dict[str, Any]):
    cats = [str(c) for c in (chart_data or {}).get("categories", ["Met","Approaches","Not Met"])]
    vals = (chart_data or {}).get("series", [{}])[0].get("values", [0,0,0])
    total = _num((chart_data or {}).get("total_students", sum(_num(v) for v in vals))) or sum(_num(v) for v in vals) or 1

    pcts = (chart_data or {}).get("status_percentages") or []
    if len(pcts) < 3 or sum(_num(x) for x in pcts) <= 1:
        pcts = [(_num(v) / total) * 100 for v in vals]

    def idx(label, fallback):
        for i, c in enumerate(cats):
            if label in c.lower():
                return i
        return fallback

    met_i = idx("met", 0)
    app_i = idx("approach", 1)
    not_i = next((i for i, c in enumerate(cats) if "not met" in c.lower()), 2)

    data = [
        ("MET", _num(vals[met_i]) if met_i < len(vals) else 0, _num(pcts[met_i]) if met_i < len(pcts) else 0, "16A34A", "ECFDF5"),
        ("APPROACHES", _num(vals[app_i]) if app_i < len(vals) else 0, _num(pcts[app_i]) if app_i < len(pcts) else 0, AMBER, "FFFBEB"),
        ("NOT MET", _num(vals[not_i]) if not_i < len(vals) else 0, _num(pcts[not_i]) if not_i < len(pcts) else 0, "E11D48", "FFF1F2"),
    ]

    # Header
    _add_text(slide, "CCMR A-F Accountability Status", 0, 27, 1333, 28, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    district = (sc.get("slide_data") or {}).get("District") or ""
    _add_text(slide, str(district), 0, 63, 1333, 17, size=13, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    label = str(sc.get("month","")).strip()
    year = str(sc.get("year_label", sc.get("year",""))).strip()
    date_label = f"{label} {year}".strip()
    if date_label:
        _add_rect(slide, 590, 88, 154, 27, NAVY, radius=True)
        _add_text(slide, date_label, 590, 96, 154, 10, size=7.8, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    _add_rect(slide, 1210, 22, 76, 31, AMBER, radius=True)
    _add_text(slide, "DATA", 1210, 31, 76, 10, size=8.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # KPI cards
    x0, y0, card_w, card_h, gap = 245, 150, 300, 94, 24
    for i, (name, count, pct, color, bg) in enumerate(data):
        x = x0 + i * (card_w + gap)
        _add_rect(slide, x, y0, card_w, card_h, bg, line=color, radius=True)
        _add_text(slide, name, x, y0 + 15, card_w, 14, size=11.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"{int(count):,}", x, y0 + 39, card_w, 23, size=21.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"{pct:.1f}%", x, y0 + 68, card_w, 12, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Progress bar
    _add_text(slide, "Progress Toward 90% Goal", 245, 270, 500, 16, size=12, color=NAVY, bold=True)
    bar_x, bar_y, bar_w, bar_h = 245, 296, 970, 29
    cursor = bar_x
    for name, count, pct, color, bg in data:
        w = max(0, bar_w * pct / 100)
        _add_rect(slide, cursor, bar_y, w, bar_h, color, radius=False)
        if w > 72:
            short = "Not Met" if name == "NOT MET" else name.title()
            _add_text(slide, f"{short} {pct:.1f}%", cursor, bar_y + 8, w, 10, size=6.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cursor += w

    goal_pct = _num((chart_data or {}).get("goal_pct", 90))
    goal_x = bar_x + bar_w * goal_pct / 100
    _add_line(slide, goal_x, bar_y - 4, goal_x, bar_y + bar_h + 8, NAVY, width=1.4)
    _add_text(slide, "GOAL", goal_x - 18, bar_y + bar_h + 12, 36, 10, size=6.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Gap box
    met_pct = data[0][2]
    needed = int(_num((chart_data or {}).get("additional_needed"))) if (chart_data or {}).get("additional_needed") is not None else int(max(0, (goal_pct/100*total) - data[0][1]) + 0.999999)
    gap_pts = _num((chart_data or {}).get("gap_pts", max(0, goal_pct - met_pct)))

    _add_rect(slide, 245, 355, 970, 82, "F3F4F6", line="F3F4F6", radius=True)
    _add_text(slide, "Gap Analysis & Path Forward", 265, 370, 450, 16, size=12.2, color=NAVY, bold=True)
    metrics = [
        ("Currently Met", f"{int(data[0][1]):,}", "16A34A"),
        ("Additional Needed", f"+{needed:,}", "E11D48"),
        ("Gap to Goal", f"{gap_pts:.1f} pts", NAVY),
        ("Total Students", f"{int(total):,}", NAVY),
    ]
    for i, (lab, val, color) in enumerate(metrics):
        x = 265 + i * 215
        _add_text(slide, lab, x, 400, 150, 11, size=7.5, color=DARK_TEXT, bold=True)
        _add_text(slide, val, x, 417, 150, 18, size=14.5, color=color, bold=True)

    # Insights / next steps boxes
    _add_rect(slide, 245, 505, 455, 92, LIGHT_BLUE, line=LIGHT_BLUE, radius=True)
    _add_text(slide, "KEY INSIGHTS", 265, 522, 260, 15, size=12.5, color=NAVY, bold=True)
    _add_text(slide, f"▶ {data[0][2]:.1f}% of students ({int(data[0][1]):,}) have met CCMR — {gap_pts:.1f} percentage points below the 90% goal.", 265, 548, 415, 18, size=8, color=DARK_TEXT)
    _add_text(slide, f"▶ {int(data[2][1]):,} students ({data[2][2]:.1f}%) have not met CCMR and need targeted support.", 265, 574, 415, 16, size=8, color=DARK_TEXT)

    _add_rect(slide, 725, 505, 490, 92, "ECFDF5", line="ECFDF5", radius=True)
    _add_text(slide, "NEXT STEPS", 745, 522, 260, 15, size=12.5, color="16A34A", bold=True)
    _add_text(slide, "▶ Prioritize students closest to meeting CCMR and assign targeted interventions before the next reporting cycle.", 745, 548, 430, 25, size=8, color=DARK_TEXT)


def _draw_ccmr_yoy_cards(slide, chart_data: Dict[str, Any], mode: str):
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", [])
    if not categories or not series:
        _add_text(slide, "No chart data available", 100, 260, 900, 50, size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return

    _add_text(slide, "YEAR-OVER-YEAR COMPARISON", 0, 155, 1333, 24, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    year_colors = [NAVY, CYAN, AMBER]
    header_colors = [NAVY, "0F766E", CYAN]
    subtitles = {
        "TSI": "CCMR TSI Status",
        "IBC": "Industry Based Certification",
        "Enrollment": "College Enrollment",
    }

    card_w, card_h = 350, 330
    start_x, gap = 75, 24
    y = 195
    max_v = 100 if mode == "percent" else max([_num(v) for s in series for v in s.get("values", [])] + [1])
    max_v = max(max_v, 1)

    for ci, metric in enumerate(categories[:3]):
        x = start_x + ci * (card_w + gap)
        _add_rect(slide, x, y, card_w, card_h, WHITE, line=GRID_GRAY, radius=True)
        _add_rect(slide, x, y, card_w, 62, header_colors[ci % len(header_colors)], line=header_colors[ci % len(header_colors)], radius=True)
        _add_text(slide, metric, x + 130, y + 16, 170, 20, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, f"({subtitles.get(metric, '')})", x + 105, y + 39, 220, 14, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        _add_oval(slide, x + 35, y + 13, 45, 45, fill=WHITE, line=WHITE, width=1)
        icon_text = "TSI" if metric == "TSI" else ("IBC" if metric == "IBC" else "HE")
        _add_text(slide, icon_text, x + 40, y + 27, 35, 12, size=9, color=header_colors[ci % len(header_colors)], bold=True, align=PP_ALIGN.CENTER)

        plot_x, plot_y, plot_w, plot_h = x + 58, y + 118, 250, 150
        _add_line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, MID_GRAY, width=0.7)
        _add_line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, MID_GRAY, width=0.7)
        ticks = [0, 20, 40, 60, 80, 100] if mode == "percent" else [0, max_v*.25, max_v*.5, max_v*.75, max_v]
        for t in ticks:
            ty = plot_y + plot_h - (_num(t) / max_v) * plot_h
            _add_line(slide, plot_x, ty, plot_x + plot_w, ty, GRID_GRAY, width=0.35)
            _add_text(slide, _fmt(t, mode), plot_x - 35, ty - 6, 30, 12, size=7, color=DARK_TEXT, align=PP_ALIGN.RIGHT)

        vals = []
        for si, s in enumerate(series[:3]):
            values = s.get("values", [])
            val = _num(values[ci]) if ci < len(values) else 0
            vals.append(val)
            bh = (val / max_v) * plot_h if max_v else 0
            bx = plot_x + 22 + si * 75
            by = plot_y + plot_h - bh
            _add_rect(slide, bx, by, 42, max(2, bh), year_colors[si % len(year_colors)], radius=True)
            _add_text(slide, _fmt(val, mode), bx - 8, by - 18, 58, 14, size=8.5, color=year_colors[si % len(year_colors)], bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, str(s.get("name", "")), bx - 6, plot_y + plot_h + 14, 55, 12, size=8, color=DARK_TEXT, align=PP_ALIGN.CENTER)

        if vals:
            diff = vals[-1] - vals[0]
            diff_text = f"{diff:+.1f} pp from {series[0].get('name','')} to {series[min(2,len(series)-1)].get('name','')}" if mode == "percent" else f"{diff:+,.0f} from {series[0].get('name','')} to {series[min(2,len(series)-1)].get('name','')}"
            _add_rect(slide, x + 35, y + 285, card_w - 70, 28, LIGHT_BLUE, line=LIGHT_BLUE, radius=True)
            _add_text(slide, f"▲ {diff_text}", x + 46, y + 293, card_w - 92, 12, size=8.5, color=header_colors[ci % len(header_colors)], bold=True, align=PP_ALIGN.CENTER)


def _draw_postsecondary_readiness_dashboard(slide, sc: Dict[str, Any]):
    sd = sc.get("slide_data", {}) or {}
    cd = sc.get("chart_data", {}) or {}
    district = sd.get("District", "District")
    month = sc.get("month", "")
    year = sc.get("year_label", "")
    footnote = _clean_source_text(sc.get("footnote", ""))
    kpis = cd.get("kpis", []) or []
    comp = cd.get("comparison", []) or []
    gaps = cd.get("opportunity_gaps", []) or []
    total = int(cd.get("total_students") or sd.get("total_students") or 0)
    comp_metric = cd.get("comparison_metric") or "Financial Aid Submitted"
    _add_rect(slide, 0, 0, 1333, 5, NAVY)
    # Dedicated header for the Postsecondary dashboard. District and date are
    # deliberately separated so the date pill never covers the district name.
    _add_text(slide, "Postsecondary Readiness", 0, 18, 1333, 32, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, district, 440, 58, 360, 18, size=12, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)
    if month or year:
        stamp = " ".join([x for x in [month, year] if x])
        _add_rect(slide, 815, 54, 100, 24, NAVY, radius=True)
        _add_text(slide, stamp, 823, 60, 84, 11, size=7.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, 1176, 22, 124, 28, "E8192C", radius=True)
    _add_text(slide, "POSTSECONDARY", 1180, 29, 116, 12, size=7.4, color=WHITE, bold=True, italic=True, align=PP_ALIGN.CENTER)
    _add_line(slide, 32, 98, 1302, 98, NAVY, width=1.4)
    _add_line(slide, 32, 100, 1302, 100, CYAN, width=1.0)
    colors = ["0057B8", CYAN, "7C3AED", "D99A00"]
    icons = ["APP", "$", "FAFSA", "4YR"]
    x0, y0, card_w, card_h, gap = 34, 112, 307, 106, 18
    for i in range(4):
        k = kpis[i] if i < len(kpis) else {}
        c = str(k.get("accent", "#"+colors[i])).replace("#", "") if k.get("accent") else colors[i]
        x = x0 + i * (card_w + gap)
        _add_rect(slide, x, y0, card_w, card_h, WHITE, line=GRID_GRAY, radius=True)
        _add_rect(slide, x, y0, 5, card_h, c)
        _add_oval(slide, x+18, y0+26, 52, 52, fill="EAF6FF", line=c, width=1.2)
        _add_text(slide, icons[i], x+23, y0+45, 42, 12, size=7.4, color=c, bold=True, align=PP_ALIGN.CENTER)
        available = bool(k.get("available")) and k.get("pct") is not None
        pct = f"{float(k.get('pct')):.1f}%" if available else "N/A"
        count_line = f"{int(k.get('count') or 0):,} of {int(k.get('total') or total):,} students" if available else "Column not available"
        delta = k.get("delta_pp")
        trend = f"+{float(delta):.1f} pp vs. prior cohort" if delta is not None else "Current cohort progress"
        _add_text(slide, k.get("label", "Metric"), x+84, y0+15, card_w-102, 20, size=10, color=NAVY, bold=True)
        _add_text(slide, pct, x+84, y0+41, card_w-102, 26, size=24, color=c, bold=True)
        _add_text(slide, count_line, x+84, y0+72, card_w-102, 11, size=7.5, color="4B5563")
        _add_text(slide, "↑ " + trend if available else trend, x+84, y0+89, card_w-102, 11, size=7.0, color="16A34A" if available else "6B7280", bold=True)
    left_x, top_y, left_w, main_h = 34, 240, 905, 302
    right_x, right_w = 955, 344
    _add_rect(slide, left_x, top_y, left_w, main_h, WHITE, line=GRID_GRAY, radius=True)
    _add_text(slide, "Campus Comparison", left_x+22, top_y+18, 300, 24, size=17, color=NAVY, bold=True)
    _add_text(slide, f"% of students with {comp_metric}", left_x+22, top_y+43, 340, 13, size=8.5, color="6B7280")
    _add_rect(slide, left_x+684, top_y+17, 190, 26, WHITE, line="0057B8", radius=True)
    _add_text(slide, f"Metric: {comp_metric}", left_x+694, top_y+24, 170, 10, size=7.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    chart_x, chart_y, chart_w, row_h = left_x+270, top_y+72, 520, 24
    for t in [0,20,40,60,80,100]:
        x = chart_x + chart_w * t / 100
        _add_line(slide, x, chart_y-8, x, chart_y+170, GRID_GRAY, width=0.4)
        _add_text(slide, f"{t}%", x-15, chart_y+178, 30, 10, size=7, color="6B7280", align=PP_ALIGN.CENTER)
    for i, r in enumerate(comp[:7], start=1):
        y = chart_y + (i-1)*row_h
        pct = float(r.get("pct") or 0)
        campus = _campus_label(r.get("campus", "Campus"))
        _add_oval(slide, left_x+24, y+4, 17, 17, fill="0057B8", line="0057B8", width=0.5)
        _add_text(slide, str(i), left_x+25, y+8, 15, 7, size=5.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, _short(campus, 34), left_x+52, y+5, 205, 12, size=8, color=NAVY, bold=True)
        _add_rect(slide, chart_x, y+5, max(2, chart_w*pct/100), 14, NAVY, radius=True)
        _add_text(slide, f"{pct:.1f}%", chart_x+chart_w*pct/100+8, y+6, 50, 10, size=7.4, color=NAVY, bold=True)
    _add_text(slide, f"Campuses ranked by % of students with {comp_metric}", left_x+24, top_y+272, 420, 10, size=6.8, color="6B7280", italic=True)
    _add_rect(slide, right_x, top_y, right_w, main_h, WHITE, line=GRID_GRAY, radius=True)
    _add_text(slide, "Opportunity Gap", right_x+22, top_y+18, 260, 22, size=16, color=NAVY, bold=True)
    _add_text(slide, "Students not yet completed", right_x+22, top_y+42, 220, 12, size=8.5, color="6B7280")
    for i, g in enumerate(gaps[:3]):
        gy = top_y+68+i*64
        c = str(g.get("accent", "#"+colors[i])).replace("#", "") if g.get("accent") else colors[i]
        _add_rect(slide, right_x+22, gy, right_w-44, 52, "F8FBFF", line=GRID_GRAY, radius=True)
        _add_oval(slide, right_x+36, gy+9, 34, 34, fill="EAF6FF", line=c, width=1)
        _add_text(slide, icons[i], right_x+40, gy+21, 26, 8, size=5.5, color=c, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, g.get("missing_label", "Missing"), right_x+84, gy+10, 145, 18, size=8.2, color=c, bold=True)
        _add_text(slide, f"of {int(g.get('total') or total):,} students", right_x+84, gy+32, 120, 9, size=6.5, color="6B7280")
        _add_text(slide, f"{int(g.get('missing_count') or 0):,}", right_x+240, gy+8, 75, 15, size=12.5, color=c, bold=True, align=PP_ALIGN.RIGHT)
        _add_text(slide, f"{float(g.get('missing_pct') or 0):.1f}%", right_x+240, gy+28, 75, 13, size=9.8, color=c, bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, "Counts may overlap across categories.", right_x+24, top_y+272, 220, 10, size=6.6, color="6B7280", italic=True)
    bx, by, bw, bh = 34, 560, 630, 108
    _add_rect(slide, bx, by, bw, bh, "EAF6FF", line="BFD7FF", radius=True)
    _add_rect(slide, bx, by, 5, bh, "0057B8")
    _add_oval(slide, bx+22, by+25, 50, 50, fill="0057B8", line="0057B8", width=1)
    _add_text(slide, "!", bx+40, by+42, 15, 12, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "KEY INSIGHTS", bx+92, by+18, 200, 14, size=12, color="0057B8", bold=True)
    bullets=[]
    if comp:
        bullets.append(f"{float(comp[0].get('pct') or 0):.1f}% at {_short(comp[0].get('campus','top campus'), 34)} for {comp_metric}.")
    if gaps:
        bullets.append(f"Largest opportunity: {int(gaps[-1].get('missing_count') or 0):,} students need support.")
    bullets.append("Use campus-level gaps to prioritize outreach.")
    for i,b in enumerate(bullets[:3]):
        _add_text(slide, "• " + b, bx+92, by+40+i*18, bw-112, 12, size=7.6, color=NAVY)
    bx2 = 680
    _add_rect(slide, bx2, by, 619, bh, "ECFDF5", line="BBF7D0", radius=True)
    _add_rect(slide, bx2, by, 5, bh, "16A34A")
    _add_oval(slide, bx2+22, by+25, 50, 50, fill="16A34A", line="16A34A", width=1)
    _add_text(slide, "->", bx2+38, by+42, 20, 12, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "NEXT STEPS", bx2+92, by+18, 200, 14, size=12, color="16A34A", bold=True)
    for i,b in enumerate(["Prioritize FAFSA and financial aid completion events by campus.", "Partner with counselors and community colleges to support completion.", "Track progress weekly and adjust support for largest gaps."]):
        _add_text(slide, "• " + b, bx2+92, by+40+i*18, 500, 12, size=7.6, color="065F46")
    _footer(slide, footnote)



def _draw_hb3_funds_dashboard(slide, sc: Dict[str, Any]):
    sd = sc.get("slide_data", {}) or {}
    cd = sc.get("chart_data", {}) or {}
    district = sd.get("District") or sd.get("Campus") or "District"
    cats = [str(c) for c in (cd.get("categories") or [])]
    series = cd.get("series") or []
    values = list(series[0].get("values", [])) if series else []
    values = [_num(v) for v in values]
    raw_statuses = cd.get("statuses") or []
    statuses = [_hb3_status_for(cats[i], raw_statuses[i] if i < len(raw_statuses) else "verified") for i in range(len(cats))]
    status_labels = {"verified":"VERIFIED", "estimate":"ESTIMATED", "projected":"PROJECTED"}
    status_colors = {"verified":NAVY, "estimate":"2563EB", "projected":CYAN}
    tag_colors = {"verified":"16A34A", "estimate":"D97706", "projected":"7C3AED"}
    total = sd.get("total_funding") or _fmt_money(sum(values))
    years_cover = sd.get("years_covered") or (f"{cats[0].replace('Class ', '')}–{cats[-1].replace('Class ', '')}" if cats else "")

    _add_rect(slide, 0, 0, 1333, 5, NAVY)
    _add_text(slide, "CCMR OUTCOMES", 24, 22, 280, 20, size=12, color=CYAN, bold=True)
    _add_text(slide, "HB3 Outcomes Bonus Funding", 24, 54, 600, 34, size=23, color=NAVY, bold=True)
    _add_text(slide, district, 24, 88, 500, 20, size=12, color=CYAN, bold=True)
    _add_line(slide, 24, 112, 1285, 112, NAVY, width=1.0)
    _add_line(slide, 24, 114, 1285, 114, CYAN, width=0.8)

    _add_text(slide, "HB3 FUNDS BY CLASS YEAR", 24, 132, 400, 18, size=12, color="6B7280", bold=True)
    x0, y0, w, h = 126, 170, 890, 415
    max_v = max(values + [1]) * 1.15
    # grid and abbreviated y-axis labels
    for k in range(5):
        t = max_v * k / 4
        y = y0 + h - (t / max_v) * h if max_v else y0+h
        _add_line(slide, x0, y, x0 + w, y, GRID_GRAY, width=0.6)
        _add_text(slide, _fmt_money(t), 42, y-7, 72, 14, size=8, color="9CA3AF", align=PP_ALIGN.RIGHT)
    _add_line(slide, x0, y0, x0, y0+h, "D1D5DB", width=0.8)
    _add_line(slide, x0, y0+h, x0+w, y0+h, "D1D5DB", width=0.8)

    n = max(1, len(cats))
    spacing = w / n
    bar_w = min(90, spacing * 0.42)
    for i, cat in enumerate(cats):
        val = values[i] if i < len(values) else 0
        st = statuses[i] if i < len(statuses) else "verified"
        bh = (val / max_v) * h if max_v else 0
        cx = x0 + spacing * i + spacing / 2
        y = y0 + h - bh
        _add_rect(slide, cx - bar_w/2, y, bar_w, max(3, bh), status_colors.get(st, NAVY), radius=True)
        if bh > 26:
            _add_text(slide, _fmt_money(val), cx - bar_w/2, y + bh/2 - 7, bar_w, 14, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, _fmt_money(val), cx - 45, y - 18, 90, 14, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, cat, cx - 70, y0+h+14, 140, 16, size=10, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, status_labels.get(st, st.upper()), cx - 70, y0+h+32, 140, 12, size=7, color="6B7280", bold=True, align=PP_ALIGN.CENTER)

    # Right side summary cards
    rx, ry, rw = 1060, 125, 238
    _add_rect(slide, rx, ry, rw, 118, NAVY, line=NAVY, radius=True)
    _add_text(slide, f"{len(cats)}-YEAR HB3 FUNDS TOTAL", rx+12, ry+16, rw-24, 12, size=7.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, str(total), rx+12, ry+38, rw-24, 32, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"{district} · Classes\n{years_cover}", rx+12, ry+76, rw-24, 32, size=12, color="BFDBFE", align=PP_ALIGN.CENTER)

    cy = ry + 130
    for i, cat in enumerate(cats[:5]):
        val = values[i] if i < len(values) else 0
        st = statuses[i] if i < len(statuses) else "verified"
        col = status_colors.get(st, NAVY)
        _add_rect(slide, rx, cy, rw, 68, "F9FAFB", line=GRID_GRAY, radius=True)
        _add_rect(slide, rx, cy, 4, 68, col)
        _add_text(slide, cat.upper(), rx+12, cy+12, 110, 14, size=10, color="6B7280", bold=True)
        _add_rect(slide, rx+170, cy+10, 58, 14, "EEF2F7", line="EEF2F7", radius=True)
        _add_text(slide, status_labels.get(st, ""), rx+172, cy+12, 54, 8, size=4.2, color=tag_colors.get(st, "6B7280"), bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, _fmt_money(val), rx+12, cy+35, 140, 20, size=16, color=col, bold=True)
        cy += 76

    _footer(slide, sc.get("footnote", ""))

def _draw_generic_chart(slide, sc: Dict[str, Any]):
    st = sc.get("slide_type", "")
    chart_data = sc.get("chart_data", {}) or {}
    mode = sc.get("mode", "count")
    if st == "ccmr_yoy_breakdown":
        _draw_ccmr_yoy_cards(slide, chart_data, mode)
    elif st == "ccmr_af_status":
        _draw_ccmr_af_status_dashboard(slide, chart_data, sc)
        return
    elif st == "ccmr_af_status_legacy":
        _draw_card_status(slide, chart_data, mode)
        _draw_grouped_columns(slide, chart_data, mode)
    elif st == "by_the_numbers":
        _draw_by_numbers(slide, chart_data, mode)
    elif st in {"tsi_status"}:
        _draw_horizontal_stacked_bars(slide, chart_data, mode)
    elif st in {"tsi_leaderboard", "ccmr_pathway_full"}:
        _draw_horizontal_bars(slide, chart_data, mode)
    elif st in {"tsi_status_trends", "postsecondary_enrollment"}:
        _draw_stacked_columns(slide, chart_data, mode)
    else:
        # For mixed/YOY slides, grouped editable columns are the safest default.
        _draw_grouped_columns(slide, chart_data, mode)


def _static_cover(slide, sc):
    sd = sc.get("slide_data", {}) or {}
    _add_rect(slide, 0, 0, 1333, 750, NAVY)
    _add_rect(slide, 0, 0, 1333, 5, RED)
    _add_rect(slide, 0, 745, 1333, 5, RED)
    # Decorative editable rings
    _add_oval(slide, 905, -150, 390, 390, fill=None, line="123A6B", width=20)
    _add_oval(slide, 970, -85, 260, 260, fill=None, line="123A6B", width=16)
    _add_oval(slide, 1130, 190, 260, 260, fill=None, line="123A6B", width=14)
    district = sd.get("District") or "District ISD × EMC"
    meeting = sd.get("meeting_type") or "Meeting Type"
    subtitle = sd.get("subtitle") or "Subtitle"
    month = sc.get("month", "")
    year = sc.get("year_label", "")
    stamp = " ".join([s for s in [month, year] if s]) or year or month
    _add_text(slide, district, 70, 295, 750, 56, size=38, color=WHITE, bold=True)
    _add_text(slide, meeting, 70, 365, 520, 36, size=20, color=CYAN, bold=True)
    _add_text(slide, subtitle, 70, 410, 520, 26, size=14, color=WHITE, italic=True)
    if stamp:
        _add_rect(slide, 70, 455, 130, 40, RED, radius=True)
        _add_text(slide, stamp, 78, 466, 114, 14, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_dark_background_logo(slide, 1088, 620, 168)


def _static_mission(slide, sc):
    """Editable mission slide matching the approved split-panel design."""
    # Right-side approved photo panel first, then left mission panel.
    if not _add_picture_asset(slide, "mission_photo_correct.png", 688, 0, 645, 750):
        _add_rect(slide, 690, 0, 643, 750, "1E3A6E")
    _add_rect(slide, 0, 0, 10, 750, RED)
    _add_rect(slide, 10, 0, 678, 750, "003DA5")
    _add_rect(slide, 10, 745, 678, 5, RED)
    _add_text(slide, "Our Mission", 60, 58, 260, 26, size=16, color=CYAN, bold=True)
    _add_rect(slide, 60, 90, 80, 4, RED)
    _add_text(slide, "Every Learner", 60, 140, 520, 56, size=36, color=WHITE, bold=True)
    _add_text(slide, "On A Path To", 60, 205, 520, 56, size=36, color=WHITE, bold=True)
    _add_text(slide, "A Living Wage", 60, 270, 520, 56, size=36, color=CYAN, bold=True)
    _add_text(slide, "While Meeting Regional Workforce Needs", 62, 360, 500, 28, size=14, color=WHITE)
    _add_dark_background_logo(slide, 72, 620, 158)


def _static_agenda(slide, sc):
    sd = sc.get("slide_data", {}) or {}
    _add_rect(slide, 0, 0, 1333, 750, NAVY)
    _add_text(slide, "TODAY'S AGENDA", 55, 55, 260, 22, size=12, color=CYAN, bold=True)
    _add_text(slide, "Presentation Overview", 55, 100, 620, 54, size=34, color=WHITE, bold=True)
    _add_line(slide, 55, 170, 1260, 170, CYAN, width=1)
    agenda = sd.get("slides_list") or []
    if not agenda:
        agenda = [{"name":"TSI Status Trends","category":"TSI"},{"name":"CCMR YOY Growth","category":"CCMR"},{"name":"Postsecondary Enrollment","category":"Postsecondary"}]
    left = agenda[:6]
    right = agenda[6:12]
    for side, items in enumerate([left, right]):
        x = 65 if side == 0 else 730
        y = 230
        last_cat = None
        for idx, item in enumerate(items, start=1 + side * 6):
            cat = item.get("category", "Other")
            if cat != last_cat:
                _add_text(slide, cat.upper(), x, y, 360, 20, size=14, color=CYAN, bold=True)
                _add_line(slide, x, y + 24, x + 550, y + 24, CYAN, width=1)
                y += 42
            _add_text(slide, str(idx), x, y, 28, 20, size=12, color=RED, bold=True)
            _add_text(slide, item.get("name", f"Slide {idx}"), x + 40, y, 470, 20, size=13, color=WHITE)
            y += 36
            last_cat = cat
    _add_dark_background_logo(slide, 66, 620, 158)


def _static_divider(slide, sc):
    title = _clean_title(sc)
    _add_rect(slide, 0, 0, 1333, 750, NAVY)
    _add_rect(slide, 0, 0, 1333, 5, RED)
    _add_rect(slide, 0, 745, 1333, 5, RED)
    _add_text(slide, title, 0, 310, 1333, 70, size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, 500, 405, 330, 5, RED)
    _add_text(slide, "ECONOMIC MOBILITY CENTER", 0, 455, 1333, 20, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


def _static_methodology(slide, sc):
    _title(slide, "Methodology", "", sc.get("month", ""), sc.get("year_label", ""), tag="")
    items = [
        "Data reflects selected districts, campuses, years, and reporting metadata.",
        "Metrics are calculated from the uploaded data and selected column mappings.",
        "Insights are generated from the same chart payload used to build each slide.",
        "Editable PowerPoint output is built with native shapes, text, and editable chart elements.",
    ]
    y = 170
    for item in items:
        _add_text(slide, "•", 95, y, 25, 24, size=18, color=CYAN, bold=True)
        _add_text(slide, item, 130, y, 930, 34, size=16, color=DARK_TEXT)
        y += 62
    _footer(slide, sc.get("footnote", ""))


def _static_outro(slide, sc):
    _add_rect(slide, 0, 0, 1333, 750, NAVY)
    _add_rect(slide, 0, 0, 1333, 5, RED)
    _add_rect(slide, 0, 745, 1333, 5, RED)
    _add_oval(slide, -80, -20, 620, 760, fill=None, line="123A6B", width=42)
    _add_oval(slide, 790, -20, 620, 760, fill=None, line="123A6B", width=42)
    _add_oval(slide, 380, 100, 570, 570, fill=None, line="2D1B4B", width=50)
    _add_text(slide, "ECONOMIC MOBILITY CENTER", 0, 288, 1333, 22, size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Thank You", 0, 330, 1333, 72, size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, 610, 420, 120, 4, RED)
    _add_text(slide, "Every Learner On A Path To A Living Wage", 0, 452, 1333, 30, size=15, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


def _draw_district_profile_dashboard(slide, chart_data: Dict[str, Any], sc: Dict[str, Any]):
    sd = sc.get("slide_data", {}) or {}
    district = sd.get("District") or sd.get("Campus") or "District"
    month = sc.get("month", "")
    year = sc.get("year_label", "")
    stamp = " ".join([x for x in [month, year] if x]).upper()
    categories = [str(c) for c in (chart_data or {}).get("categories", [])]
    series = (chart_data or {}).get("series", []) or []
    bar_colors = ["93C5FD", "1D4ED8", NAVY, CYAN]

    _add_rect(slide, 0, 0, 1333, 5, NAVY)
    _add_text(slide, "Economic Mobility Center District Profile", 26, 20, 760, 34, size=20, color=NAVY, bold=True)
    _add_text(slide, str(district).upper(), 26, 58, 760, 20, size=12, color=CYAN, bold=True)
    if stamp:
        _add_text(slide, stamp, 1090, 28, 210, 20, size=10, color="6B7280", align=PP_ALIGN.RIGHT)
    _add_line(slide, 26, 88, 1280, 88, NAVY, width=1.0)
    _add_line(slide, 26, 90, 1280, 90, CYAN, width=0.7)

    lx = 30
    for i, cat in enumerate(categories[:6]):
        _add_rect(slide, lx, 105, 14, 14, bar_colors[i % len(bar_colors)], radius=True)
        _add_text(slide, cat, lx + 20, 103, 95, 18, size=8, color=DARK_TEXT, bold=True)
        lx += 118

    x0, y0, w, h = 28, 138, 1276, 532
    _add_rect(slide, x0, y0, w, h, WHITE, line=MID_GRAY, radius=True)
    if not series:
        _add_text(slide, "No District Profile data available", x0, y0 + 220, w, 40, size=20, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return

    def metric_scale(values):
        # District Profile uses a true 0-100% axis across all metric panels.
        # Readability is handled by larger labels, not by changing the scale.
        return 100.0

    tile_w = w / max(1, len(series))
    plot_h = 420
    base_y = y0 + h - 82
    top_y = base_y - plot_h

    for mi, metric in enumerate(series):
        tx = x0 + mi * tile_w
        if mi > 0:
            _add_line(slide, tx, y0, tx, y0 + h, MID_GRAY, width=1.0)
        label = str(metric.get("name", "Metric"))
        _add_text(slide, label, tx + 8, y0 + 14, tile_w - 16, 30, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        values = metric.get("values", []) or []
        available = bool(metric.get("available", any(v is not None for v in values)))
        local_max = metric_scale(values)
        _add_line(slide, tx + 18, base_y, tx + tile_w - 18, base_y, MID_GRAY, width=1.0)
        if not available:
            _add_text(slide, "N/A", tx + 8, y0 + 225, tile_w - 16, 26, size=18, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, "Column not available", tx + 8, y0 + 253, tile_w - 16, 16, size=8, color="6B7280", align=PP_ALIGN.CENTER)
            continue
        n = max(1, len(categories))
        gap = (tile_w - 42) / n
        bar_w = min(50, gap * 0.84)
        for ci, cat in enumerate(categories):
            cx = tx + 25 + gap * ci + gap / 2
            val = None
            if ci < len(values):
                try:
                    val = float(values[ci]) if values[ci] is not None else None
                except Exception:
                    val = None
            if val is None:
                _add_text(slide, "N/A", cx - 22, base_y - 24, 44, 14, size=7.5, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER)
            else:
                bh = (max(0, val) / local_max) * plot_h if local_max else 0
                bh = max(10, bh) if val > 0 else max(2, bh)
                _add_rect(slide, cx - bar_w / 2, base_y - bh, bar_w, bh, bar_colors[ci % len(bar_colors)], radius=True)
                # Put all data labels at the top of the bars so the PPTX values stay readable.
                label_y = max(top_y + 8, base_y - bh - 18)
                _add_text(slide, f"{val:.1f}%", cx - 31, label_y + 4, 62, 14, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, str(cat), cx - 32, base_y + 14, 64, 18, size=8.5, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)




def _data_slide(slide, sc):
    if sc.get("slide_type", "") == "hb3_funds":
        _draw_hb3_funds_dashboard(slide, sc)
        return

    if sc.get("slide_type", "") == "postsecondary_enrollment" and ((sc.get("chart_data", {}) or {}).get("dashboard_type") == "postsecondary_readiness" or (sc.get("layout", "") == "postsecondary_readiness")):
        _draw_postsecondary_readiness_dashboard(slide, sc)
        return

    if sc.get("slide_type", "") == "district_profile":
        _draw_district_profile_dashboard(slide, sc.get("chart_data", {}) or {}, sc)
        _footer(slide, sc.get("footnote", ""))
        return

    if sc.get("slide_type", "") == "ccmr_pathway":
        _draw_ccmr_pathway_dashboard(slide, sc.get("chart_data", {}) or {}, sc)
        _footer(slide, sc.get("footnote", ""))
        return

    if sc.get("slide_type", "") == "ccmr_af_status":
        _draw_ccmr_af_status_dashboard(slide, sc.get("chart_data", {}) or {}, sc)
        _footer(slide, sc.get("footnote", ""))
        return

    sd = sc.get("slide_data", {}) or {}
    title = sd.get("Title") or _clean_title(sc)
    district = sd.get("District") or sd.get("Campus") or ""
    _title(slide, title, district, sc.get("month", ""), sc.get("year_label", ""), tag="DATA")
    _draw_generic_chart(slide, sc)
    _add_insight_boxes(slide, sc.get("insights", []))
    _footer(slide, sc.get("footnote", ""))



def _build_slide(prs: Presentation, sc: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_type = sc.get("slide_type", "")
    if slide_type == "cover":
        _static_cover(slide, sc)
    elif slide_type == "mission":
        _static_mission(slide, sc)
    elif slide_type == "agenda":
        _static_agenda(slide, sc)
    elif slide_type == "section_divider":
        _static_divider(slide, sc)
    elif slide_type == "methodology":
        _static_methodology(slide, sc)
    elif slide_type == "outro":
        _static_outro(slide, sc)
    else:
        _data_slide(slide, sc)
    return slide


def generate_pptx_file(slides_config: Sequence[Dict[str, Any]], out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for sc in slides_config or []:
        _build_slide(prs, sc)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path
