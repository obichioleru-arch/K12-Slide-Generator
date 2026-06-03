"""
District Slide Tool — Backend v6
Templates are now clean (no floating label text boxes).
All charts use built-in data labels.
"""

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import pandas as pd
import shutil, os, json, traceback
from datetime import datetime

from calculations.tsi import (
    calculate_tsi_status_trends,
    calculate_tsi_status,
    calculate_tsi_leaderboard,
    REQUIRED_FIELDS as TSI_FIELDS,
)
from calculations.ccmr import (
    calculate_ccmr_yoy_breakdown,
    REQUIRED_FIELDS as CCMR_FIELDS,
)
from calculations.postsecondary import (
    calculate_postsecondary_enrollment,
    REQUIRED_FIELDS as POST_FIELDS,
)

app = FastAPI(title="District Slide Tool", version="6.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173","http://localhost:5174","http://127.0.0.1:5173","http://127.0.0.1:5174"],
    allow_credentials=True, allow_methods=["*"], allow_headers=["*"],
)
os.makedirs("uploads", exist_ok=True)
os.makedirs("outputs", exist_ok=True)

# ── Registry ──────────────────────────────────────────────────────────────────
CATEGORY_MENU = {
    "TSI": [
        {"slide_name": "TSI Status Trends", "slide_type": "tsi_status_trends"},
        {"slide_name": "TSI Status",         "slide_type": "tsi_status"},
        {"slide_name": "TSI Leaderboard",    "slide_type": "tsi_leaderboard"},
    ],
    "CCMR A-F Status": [
        {"slide_name": "CCMR YOY Breakdown", "slide_type": "ccmr_yoy_breakdown"},
    ],
    "Postsecondary Enrollment": [
        {"slide_name": "Postsecondary Enrollment", "slide_type": "postsecondary_enrollment"},
    ],
    "HB3": [], "District Profile": [], "Ad Hoc": [],
}

ALL_REQUIRED_FIELDS = {}
ALL_REQUIRED_FIELDS.update(TSI_FIELDS)
ALL_REQUIRED_FIELDS.update(CCMR_FIELDS)
ALL_REQUIRED_FIELDS.update(POST_FIELDS)

MANUAL_TEXT_FIELDS = [
    {"key": "District", "label": "District / Organization Name",
     "description": "Top-left label on the slide (e.g. 'Grand Prairie ISD')"},
    {"key": "Campus",   "label": "Campus / Subtitle",
     "description": "Subtitle under the slide title (leave blank if not needed)"},
    {"key": "month",    "label": "Month (Leaderboard only)",
     "description": "Month shown in subtitle (e.g. 'October')"},
    {"key": "year",     "label": "Report Year (Leaderboard only)",
     "description": "Year shown in subtitle (e.g. '2025')"},
]

SLIDE_REGISTRY = {
    "tsi_status_trends": {
        "template": "templates/tsi_status_template.pptx",
        "output_prefix": "tsi_status_trends",
        "calculator": calculate_tsi_status_trends,
        "layout_profile": "tsi_status_trends",
        "supports_modes": ["count", "percent"],
    },
    "tsi_status": {
        "template": "templates/tsi_status_slide_template.pptx",
        "output_prefix": "tsi_status",
        "calculator": calculate_tsi_status,
        "layout_profile": "tsi_status",
        "supports_modes": ["count", "percent"],
    },
    "tsi_leaderboard": {
        "template": "templates/tsi_leaderboard_template.pptx",
        "output_prefix": "tsi_leaderboard",
        "calculator": calculate_tsi_leaderboard,
        "layout_profile": "tsi_leaderboard",
        "supports_modes": ["percent"],
    },
    "ccmr_yoy_breakdown": {
        "template": "templates/ccmr_yoy_breakdown_template.pptx",
        "output_prefix": "ccmr_yoy_breakdown",
        "calculator": calculate_ccmr_yoy_breakdown,
        "layout_profile": "ccmr_yoy_breakdown",
        "supports_modes": ["count", "percent"],
    },
    "postsecondary_enrollment": {
        "template": "templates/postsecondary_enrollment_template.pptx",
        "output_prefix": "postsecondary_enrollment",
        "calculator": calculate_postsecondary_enrollment,
        "layout_profile": "postsecondary_enrollment",
        "supports_modes": ["percent"],
    },
}

INST_CANDIDATES = [
    "Primary Educational Institution","District","District Name",
    "Campus","Campus Name","School","School Name","High School",
]


# ── Endpoints ─────────────────────────────────────────────────────────────────
@app.get("/")
def home(): return {"message": "District Slide Tool v6"}

@app.get("/health")
def health():
    return {
        "status": "ok",
        "templates_found": [
            {"slide_type": k, "template": v["template"], "exists": os.path.isfile(v["template"])}
            for k, v in SLIDE_REGISTRY.items()
        ],
        "registered_slides": list(SLIDE_REGISTRY.keys()),
    }

@app.get("/category-menu")
def get_category_menu(): return {"category_menu": CATEGORY_MENU}

@app.get("/categories")
def get_categories(): return {"categories": list(CATEGORY_MENU.keys())}

@app.get("/slide-types")
def get_slide_types(): return {"available_slides": list(SLIDE_REGISTRY.keys())}

@app.get("/slide-fields/{slide_type}")
def get_slide_fields(slide_type: str):
    if slide_type not in SLIDE_REGISTRY:
        raise HTTPException(404, f"Unknown: {slide_type}")
    return {
        "slide_type": slide_type,
        "fields": ALL_REQUIRED_FIELDS.get(slide_type, []),
        "manual_text_fields": MANUAL_TEXT_FIELDS,
        "supports_modes": SLIDE_REGISTRY[slide_type].get("supports_modes", ["count"]),
    }


@app.post("/inspect-file")
async def inspect_file(slide_type: str = Form(...), file: UploadFile = File(...)):
    if slide_type not in SLIDE_REGISTRY:
        raise HTTPException(400, f"Unknown slide type: {slide_type}")
    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()
    if ext not in (".xlsx", ".xls", ".csv"):
        raise HTTPException(400, f"Unsupported file type '{ext}'")

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    upload_path = os.path.join("uploads", f"{ts}_{filename}")
    with open(upload_path, "wb") as buf:
        shutil.copyfileobj(file.file, buf)

    sheets_info = []
    if ext == ".csv":
        df = pd.read_csv(upload_path)
        sheets_info.append(_inspect_sheet(df, "Sheet1", slide_type))
    else:
        xl = pd.ExcelFile(upload_path)
        for name in xl.sheet_names:
            try:
                df = pd.read_excel(xl, sheet_name=name)
                sheets_info.append(_inspect_sheet(df, name, slide_type))
            except Exception as e:
                sheets_info.append({
                    "sheet_name": name, "columns": [], "row_count": 0,
                    "institutions": [], "years": [], "detected_fields": [],
                    "all_required": False, "usable": False, "error": str(e),
                })

    # Aggregate all institutions across ALL usable sheets for cross-tab/cross-district use
    all_institutions = {}  # inst_name -> list of sheet_names it appears in
    for si in sheets_info:
        if si.get("usable"):
            for inst in si.get("institutions", []):
                all_institutions.setdefault(inst, []).append(si["sheet_name"])

    return {
        "upload_path": upload_path,
        "file_ext": ext,
        "sheets": sheets_info,
        "all_institutions": [
            {"name": k, "sheets": v} for k, v in sorted(all_institutions.items())
        ],
        "manual_text_fields": MANUAL_TEXT_FIELDS,
    }


def _inspect_sheet(df: pd.DataFrame, sheet_name: str, slide_type: str) -> dict:
    from calculations.tsi import _find_column as _fc
    columns    = list(df.columns)
    field_defs = ALL_REQUIRED_FIELDS.get(slide_type, [])
    detected_fields = []
    hard_missing = 0
    for field in field_defs:
        detected = _fc(df, field["candidates"])
        optional = field.get("optional", False)
        if not detected and not optional:
            hard_missing += 1
        detected_fields.append({
            "key": field["key"], "label": field["label"],
            "description": field["description"],
            "optional": optional,
            "candidates": field["candidates"],
            "detected": detected,
        })
    from calculations.tsi import _find_column as _fc
    inst_col = _fc(df, INST_CANDIDATES)
    institutions = sorted(df[inst_col].dropna().astype(str).unique().tolist()) if inst_col else []
    year_col = _fc(df, ["Class","Year","Graduation Year","Cohort","School Year"])
    years = []
    if year_col:
        raw = pd.to_numeric(df[year_col], errors="coerce").dropna()
        years = sorted(int(y) for y in raw.unique() if 2018 <= y <= 2030)
    usable = len(columns) > 2 and hard_missing < len([f for f in field_defs if not f.get("optional")])
    return {
        "sheet_name": sheet_name, "columns": columns,
        "row_count": len(df), "institutions": institutions,
        "institution_col": inst_col, "years": years,
        "detected_fields": detected_fields,
        "all_required": hard_missing == 0, "usable": usable,
    }


@app.post("/detect-columns")
async def detect_columns(
    slide_type:            str = Form(...),
    upload_path:           str = Form(...),
    sheet_name:            str = Form(None),
    selected_institutions: str = Form("[]"),
):
    if not os.path.isfile(upload_path):
        raise HTTPException(400, "File not found.")
    df = _load_df(upload_path, sheet_name)
    try:
        institutions = json.loads(selected_institutions)
    except Exception:
        institutions = []
    if institutions:
        from calculations.tsi import _find_column as _fc
        ic = _fc(df, INST_CANDIDATES)
        if ic:
            df = df[df[ic].astype(str).isin(institutions)]

    from calculations.tsi import _find_column as _fc
    field_defs = ALL_REQUIRED_FIELDS.get(slide_type, [])
    detection  = []
    for field in field_defs:
        detected = _fc(df, field["candidates"])
        detection.append({
            "key": field["key"], "label": field["label"],
            "description": field["description"],
            "optional": field.get("optional", False),
            "detected": detected, "candidates": field["candidates"],
        })
    hard_missing = [f for f in detection if not f["optional"] and not f["detected"]]
    return {
        "slide_type": slide_type,
        "file_columns": list(df.columns),
        "upload_path": upload_path,
        "sheet_name": sheet_name,
        "fields": detection,
        "hard_missing": [f["label"] for f in hard_missing],
        "can_proceed": len(hard_missing) == 0,
    }


@app.post("/preview-slide")
async def preview_slide(
    slide_type:            str = Form(...),
    upload_path:           str = Form(...),
    sheet_name:            str = Form(None),
    selected_institutions: str = Form("[]"),
    overrides:             str = Form("{}"),
    manual_text:           str = Form("{}"),
    mode:                  str = Form("count"),
):
    if not os.path.isfile(upload_path):
        raise HTTPException(400, "File not found.")
    df = _load_filtered_df(upload_path, sheet_name, selected_institutions)
    try:
        override_map = json.loads(overrides)
        manual_map   = json.loads(manual_text)
    except Exception:
        override_map = {}; manual_map = {}
    try:
        result = SLIDE_REGISTRY[slide_type]["calculator"](df, overrides=override_map, mode=mode)
    except ValueError as e:
        raise HTTPException(400, {"error": str(e)})
    except Exception as e:
        print(traceback.format_exc())
        raise HTTPException(500, {"error": f"Calculation error: {e}"})
    slide_data = result.get("slide_data", {})
    for k, v in manual_map.items():
        if v and str(v).strip():
            slide_data[k] = str(v).strip()
    return {"slide_data": slide_data, "chart_data": result.get("chart_data"), "mode": mode}


@app.post("/generate-slide")
async def generate_slide(
    slide_type:            str        = Form(...),
    file:                  UploadFile = File(None),
    upload_path:           str        = Form(None),
    sheet_name:            str        = Form(None),
    selected_institutions: str        = Form("[]"),
    overrides:             str        = Form("{}"),
    manual_text:           str        = Form("{}"),
    mode:                  str        = Form("count"),
    preview_slide_data:    str        = Form(None),
    preview_chart_data:    str        = Form(None),
):
    if slide_type not in SLIDE_REGISTRY:
        raise HTTPException(400, {"error": "Unknown slide type"})
    config = SLIDE_REGISTRY[slide_type]
    if not os.path.isfile(config["template"]):
        raise HTTPException(500, {"error": f"Template not found: {config['template']}"})

    if upload_path and os.path.isfile(upload_path):
        final_path = upload_path
    elif file:
        filename = file.filename or ""
        ext = os.path.splitext(filename)[1].lower()
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        final_path = os.path.join("uploads", f"{ts}_{filename}")
        with open(final_path, "wb") as buf:
            shutil.copyfileobj(file.file, buf)
    else:
        raise HTTPException(400, {"error": "No file provided."})

    try:
        override_map = json.loads(overrides)
        manual_map   = json.loads(manual_text)
    except Exception:
        override_map = {}; manual_map = {}

    if preview_slide_data and preview_chart_data:
        try:
            slide_data    = json.loads(preview_slide_data)
            chart_payload = json.loads(preview_chart_data)
        except Exception:
            slide_data = {}; chart_payload = None
    else:
        df = _load_filtered_df(final_path, sheet_name, selected_institutions)
        try:
            result = config["calculator"](df, overrides=override_map, mode=mode)
        except ValueError as e:
            raise HTTPException(400, {"error": str(e)})
        except Exception as e:
            print(traceback.format_exc())
            raise HTTPException(500, {"error": f"Calculation error: {e}"})
        slide_data    = result.get("slide_data", {})
        chart_payload = result.get("chart_data")

    for k, v in manual_map.items():
        if v and str(v).strip():
            slide_data[k] = str(v).strip()

    try:
        prs = Presentation(config["template"])
    except Exception as e:
        raise HTTPException(500, {"error": f"Could not open template: {e}"})

    replace_named_placeholders(prs, slide_data)

    if chart_payload:
        try:
            update_chart(prs, chart_payload, config["layout_profile"], mode)
        except Exception as e:
            print(traceback.format_exc())
            raise HTTPException(500, {"error": f"Chart update failed: {e}"})

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_name = f"{config['output_prefix']}_{mode}_{ts}.pptx"
    out_path = os.path.join("outputs", out_name)
    prs.save(out_path)

    return FileResponse(
        out_path, filename=out_name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ── File helpers ──────────────────────────────────────────────────────────────
def _load_df(path: str, sheet_name=None) -> pd.DataFrame:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":   return pd.read_csv(path)
    elif ext == ".xls": return pd.read_excel(path, engine="xlrd", sheet_name=sheet_name or 0)
    else:               return pd.read_excel(path, sheet_name=sheet_name or 0)


def _load_filtered_df(path: str, sheet_name, sel_json: str) -> pd.DataFrame:
    df = _load_df(path, sheet_name)
    try:
        institutions = json.loads(sel_json) if sel_json else []
    except Exception:
        institutions = []
    if institutions:
        from calculations.tsi import _find_column as _fc
        ic = _fc(df, INST_CANDIDATES)
        if ic:
            df = df[df[ic].astype(str).isin(institutions)]
    return df


# ── Placeholder replacement ───────────────────────────────────────────────────
def replace_named_placeholders(prs, slide_data: dict):
    """Replace {{Key}} tokens in all text shapes. Style {{District}} specially."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            raw = shape.text
            if "{{" not in raw:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for key, val in slide_data.items():
                        run.text = run.text.replace(f"{{{{{key}}}}}", str(val))

            # Style District label
            if "{{District}}" in raw:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size  = Pt(13)
                        run.font.bold  = True
                        run.font.name  = "Calibri"
                        run.font.color.rgb = RGBColor(0, 176, 240)

    # Remove any shapes that still contain unreplaced {{ }} tokens
    for slide in prs.slides:
        to_rm = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                raw = shape.text.strip()
                if raw.startswith("{{") and raw.endswith("}}") and raw[2:-2] not in slide_data:
                    to_rm.append(shape)
        for s in to_rm:
            try: s._element.getparent().remove(s._element)
            except Exception: pass


# ── Chart update & formatting ─────────────────────────────────────────────────
def update_chart(prs, chart_payload: dict, layout_profile: str, mode: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            cd = CategoryChartData()
            cd.categories = chart_payload["categories"]
            for s in chart_payload["series"]:
                cd.add_series(s["name"], tuple(s["values"]))
            shape.chart.replace_data(cd)
            _fmt_chart(shape, shape.chart, layout_profile, mode)
            return


def _fmt_chart(shape, chart, layout_profile: str, mode: str):
    """
    Apply consistent formatting to every chart.
    ALL charts get data labels — no exceptions.
    """
    fmt_pct  = '0.0"%"'
    fmt_num  = "0"
    fmt      = fmt_pct if mode == "percent" else fmt_num

    # ── shared axis cleanup ────────────────────────────────────────────────
    try:
        chart.category_axis.has_title = False
        chart.category_axis.tick_labels.font.size = Pt(11)
    except Exception: pass
    try:
        chart.value_axis.has_title = False
    except Exception: pass

    # ── per-profile settings ───────────────────────────────────────────────
    if layout_profile == "tsi_leaderboard":
        # Horizontal bar — keep template position
        try:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 100
            chart.value_axis.major_unit    = 10
        except Exception: pass
        colors = [RGBColor(0, 50, 145)]
        label_pos = 3  # OUTSIDE_END
        label_color = RGBColor(31, 41, 51)
        label_size  = Pt(11)
        for idx, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[idx % len(colors)]
            s.has_data_labels = True
            s.data_labels.position       = label_pos
            s.data_labels.show_value     = True
            s.data_labels.number_format  = fmt
            s.data_labels.font.size      = label_size
            s.data_labels.font.bold      = True
            s.data_labels.font.color.rgb = label_color

    elif layout_profile in ("tsi_status_trends", "tsi_status"):
        # Stacked column — labels inside each segment
        colors = [RGBColor(0, 50, 145), RGBColor(0, 176, 240), RGBColor(198, 40, 40)]
        for idx, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[idx % len(colors)]
            s.has_data_labels = True
            s.data_labels.position      = 1  # CENTER
            s.data_labels.show_value    = True
            s.data_labels.number_format = fmt
            s.data_labels.font.size     = Pt(9)
            s.data_labels.font.bold     = True
            s.data_labels.font.color.rgb = RGBColor(255, 255, 255)

    elif layout_profile == "ccmr_yoy_breakdown":
        # Grouped column — labels outside end
        try:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 100 if mode == "percent" else None
            if mode == "percent":
                chart.value_axis.major_unit = 10
        except Exception: pass
        colors = [RGBColor(255, 192, 0), RGBColor(0, 176, 240), RGBColor(0, 50, 145)]
        for idx, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[idx % len(colors)]
            s.has_data_labels = True
            s.data_labels.position      = 3  # OUTSIDE_END
            s.data_labels.show_value    = True
            s.data_labels.number_format = fmt
            s.data_labels.font.size     = Pt(9)
            s.data_labels.font.bold     = True
            s.data_labels.font.color.rgb = RGBColor(31, 41, 51)

    elif layout_profile == "postsecondary_enrollment":
        # Stacked column — labels inside, hide zeros
        try:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 120
            chart.value_axis.major_unit    = 20
        except Exception: pass
        colors = [RGBColor(0,84,114), RGBColor(0,176,240), RGBColor(255,192,0), RGBColor(218,150,216)]
        for idx, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[idx % len(colors)]
            s.has_data_labels = True
            s.data_labels.position      = 1  # CENTER
            s.data_labels.show_value    = True
            s.data_labels.number_format = fmt_pct  # always percent for postsecondary
            s.data_labels.font.size     = Pt(9)
            s.data_labels.font.bold     = True
            s.data_labels.font.color.rgb = RGBColor(255, 255, 255)
        try:
            chart.plots[0].overlap = 100
        except Exception: pass
